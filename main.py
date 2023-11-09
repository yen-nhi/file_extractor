import fitz
from PIL import Image
import io
from typing import Tuple
from pptx import Presentation
from deep_translator import GoogleTranslator
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import typer

# Create a python cli app to run script by command
app = typer.Typer()


def transform_RGB_from_int(RGBint: int) -> Tuple:
    """
    :param RGBint: an interger represent for a color
    :return: a tuple with RGB (red, green, blue) in range 0-1
    """
    blue = RGBint & 255
    green = (RGBint >> 8) & 255
    red = (RGBint >> 16) & 255
    return round(red / 255, 3), round(green / 255, 3), round(blue / 255, 3)


@app.command('pdf_extract')
def pdf_extract(file_path: str = 'pdf_mock_file.pdf') -> None:
    """
    Extract all images and text, store images in folder images, uppercase text and export to another pdf file
    :param file_path: the path to the file
    :return:
    """
    with fitz.open(file_path) as doc:
        output = fitz.open()
        new_page = output.new_page(-1, width=595, height=842)  # A4 dimension
        x, y = 50, 50  # position to write the text (x: from left, y: from top)
        image_count = 0

        for page in doc:
            elements = page.get_text('dict')
            for block in elements.get('blocks', []):

                # Type 1 => image
                if block.get('type') == 1:
                    image = Image.open(io.BytesIO(block.get('image')))
                    image_count += 1
                    image.save(open(f"images/pdf_image_{image_count}.{block.get('ext')}", "wb"))

                # Type 0 => text
                elif block.get('type') == 0:
                    lines = block.get('lines')
                    for line in lines:
                        for span in line.get('spans'):
                            y += 15
                            if y > 792:
                                new_page = output.new_page(-1, width=595, height=842)
                                y = 50
                            point = fitz.Point(x, y)

                            size = span.get('size')  # Extract font size
                            font = span.get('font', '').split('-')
                            font_name = font[0]  # Extract font name
                            font_style = font[1] if len(font) > 1 else None  # Extract style (italic, bold...)
                            color = transform_RGB_from_int(span.get('color'))  # Extract color
                            text = span.get('text')  # Extract text string

                            # print(
                            #     f'font size = {size}, font name = {font_name}, font style = {font_style}, '
                            #     f'color = {color},text = {text}')

                            new_page.insert_text(point=point, text=text.upper(), fontsize=size - 2, color=color)

        output.save('pdf_text.pdf')


@app.command('pptx_translate')
def pptx_translate(file_path: str = 'Networking.pptx', source: str='auto', target: str='en'):
    """
    Translate PowerPoint (.pptx) file and add target language texts to the existing file
    :param file_path:
    :param source:
    :param target:
    :return:
    """
    translator = GoogleTranslator(source=source, target=target)
    pptx = Presentation(file_path)
    for slide in pptx.slides:
        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pic = shape.image
                image = Image.open(io.BytesIO(pic.blob))
                image.save(open(f"images/pptx_image_{idx}.{pic.ext}", "wb"))
            elif shape.has_text_frame:
                text_frame = shape.text_frame
                for p in text_frame.paragraphs:
                    if not p.text:
                        continue
                    translated = translator.translate(p.text)
                    p.add_line_break()
                    new_line = p.add_run()
                    new_line.text = translated
                    new_line.font.size = Pt(12)

    pptx.save('pptx_output.pptx')


if __name__ == '__main__':
    pdf_extract()
    pptx_translate(source='en', target='vi')
